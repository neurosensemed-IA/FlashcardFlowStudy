import streamlit as st
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
import pandas as pd
import io
import google.generativeai as genai
import json
import plotly.graph_objects as go # Importar Plotly

# --- Configuración de la Página ---
st.set_page_config(
    page_title="Med-Flash AI",
    page_icon="🧬",
    layout="wide",
    initial_sidebar_state="expanded",
)

# --- ESTILOS CSS ---
st.markdown("""
<style>
    /* Paleta de colores */
    :root {
        --primary-color: #F5A6C1; /* Rosa Principal */
        --secondary-color: #E0E0E0; /* Gris Claro */
        --text-color: #4A4A4A; /* Gris Oscuro */
        --bg-color: #FFFFFF; /* Blanco */
        --dark-bg: #1E1E1E; /* Fondo oscuro opcional */
        --dark-text: #F0F0F0; /* Texto claro opcional */
    }

    /* Estilo para tema oscuro (preferido por Streamlit) */
    body {
        background-color: var(--dark-bg);
        color: var(--dark-text);
    }
    
    /* Contenedor principal */
    .stApp {
        background-color: var(--dark-bg);
    }

    /* Barra lateral */
    [data-testid="stSidebar"] {
        background-color: #2F2F2F;
        border-right: 2px solid var(--primary-color);
    }
    [data-testid="stSidebar"] .stButton button {
        background-color: transparent;
        color: var(--dark-text);
        border: 2px solid var(--primary-color);
        border-radius: 12px;
        width: 100%;
        margin-bottom: 10px;
    }
    [data-testid="stSidebar"] .stButton button:hover {
        background-color: var(--primary-color);
        color: var(--text-color);
        border-color: var(--primary-color);
    }
    [data-testid="stSidebar"] .stRadio > label {
        color: var(--dark-text) !important;
    }

    /* Botones principales */
    .stButton > button {
        background-color: var(--primary-color);
        color: var(--text-color);
        font-weight: bold;
        border-radius: 12px;
        padding: 10px 20px;
        border: none;
    }
    .stButton > button:hover {
        background-color: #F7BACF;
        color: var(--text-color);
    }

    /* Estilo de Tarjetas (Flashcards) */
    .flashcard {
        background-color: #2F2F2F; /* Fondo de tarjeta oscuro */
        border-radius: 12px;
        padding: 24px;
        margin-top: 20px;
        margin-bottom: 20px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.4);
        border: 1px solid #4A4A4A;
        color: var(--dark-text); /* Texto dentro de la tarjeta */
    }
    .flashcard h5 {
        color: var(--primary-color); /* Título de la pregunta en rosa */
        margin-bottom: 15px;
        font-size: 1.25rem;
    }

    /* Cajas de Alerta (Info, Success, Error) */
    [data-testid="stAlert"] {
        border-radius: 12px;
    }
    [data-testid="stAlert"] [data-testid="stMarkdownContainer"] p {
        color: #000; /* Texto oscuro para mejor legibilidad en alertas */
    }

    /* Contenedores de Feedback (más coloridos) */
    .feedback-correct {
        background-color: #2F2F2F;
        border: 2px solid #28a745; /* Verde */
        border-radius: 12px;
        padding: 16px;
        margin-top: 10px;
        color: #F0F0F0;
    }
    .feedback-incorrect {
        background-color: #2F2F2F;
        border: 2px solid #dc3545; /* Rojo */
        border-radius: 12px;
        padding: 16px;
        margin-top: 10px;
        color: #F0F0F0;
    }
    .feedback-explanation {
        background-color: #2F2F2F;
        border: 2px solid #17a2b8; /* Azul info */
        border-radius: 12px;
        padding: 16px;
        margin-top: 10px;
        color: #F0F0F0;
    }

    /* Contenedor de "Doodle" (Ahora con SVG) */
    .doodle-container {
        width: 100%;
        height: 150px;
        background-color: var(--primary-color);
        border-radius: 12px;
        display: flex;
        align-items: center;
        justify-content: center;
        margin-bottom: 20px;
        padding: 10px;
    }
    .doodle-container svg {
        max-width: 80%;
        max-height: 80%;
        fill: var(--text-color); /* Color de relleno para el SVG */
    }
</style>
""", unsafe_allow_html=True)

# --- Funciones de Extracción ---
def extraer_texto_pdf(file_stream):
    try:
        doc = fitz.open(stream=file_stream.read(), filetype="pdf")
        texto = ""
        for page in doc:
            texto += page.get_text()
        doc.close()
        return texto
    except Exception as e:
        return f"Error al procesar PDF: {e}"

def extraer_texto_pptx(file_stream):
    try:
        prs = Presentation(file_stream)
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
        return texto
    except Exception as e:
        return f"Error al procesar PPTX: {e}"

# --- Estado de Sesión ---
if 'page' not in st.session_state:
    st.session_state.page = "Cargar Contenido"
if 'api_key' not in st.session_state:
    st.session_state.api_key = ""
if 'user_name' not in st.session_state:
    st.session_state.user_name = ""
if 'extracted_content' not in st.session_state:
    st.session_state.extracted_content = None
if 'flashcard_library' not in st.session_state:
    st.session_state.flashcard_library = {} # Biblioteca para guardar mazos
if 'current_exam' not in st.session_state:
    st.session_state.current_exam = None
if 'current_question_index' not in st.session_state:
    st.session_state.current_question_index = 0
if 'user_answer' not in st.session_state:
    st.session_state.user_answer = None
if 'show_explanation' not in st.session_state:
    st.session_state.show_explanation = False
if 'exam_results' not in st.session_state:
    st.session_state.exam_results = []

# --- Funciones de Callback ---
def go_to_next_question():
    """Avanza a la siguiente pregunta y resetea el estado."""
    st.session_state.current_question_index += 1
    st.session_state.user_answer = None
    st.session_state.show_explanation = False

def restart_exam():
    """Reinicia el examen limpiando el estado."""
    st.session_state.current_exam = None
    st.session_state.current_question_index = 0
    st.session_state.user_answer = None
    st.session_state.show_explanation = False
    st.session_state.exam_results = []

# --- BARRA LATERAL (Navegación) ---
with st.sidebar:
    st.title("Med-Flash AI 🧬")
    st.markdown("Tu asistente de estudio médico con IA.")
    
    # SVG de Flashcard Médica (Corazón y Cerebro)
    st.markdown(f"""
    <div class="doodle-container">
        <svg viewBox="0 0 24 24" fill="currentColor">
            <path d="M19 3H5C3.89543 3 3 3.89543 3 5V19C3 20.1046 3.89543 21 5 21H19C20.1046 21 21 20.1046 21 19V5C21 3.89543 20.1046 3 19 3ZM19 5V19H5V5H19Z"></path>
            <path d="M17 7H7V17H17V7Z" fill="var(--primary-color)"></path>
            <path d="M12 8C10.6667 8 9.33333 9.33333 8 10C9.33333 10.6667 10.6667 12 12 12C13.3333 12 14.6667 10.6667 16 10C14.6667 9.33333 13.3333 8 12 8Z" fill="var(--text-color)"></path>
            <path d="M12 13C10.6667 13 9.33333 14.3333 8 15C9.33333 15.6667 10.6667 17 12 17C13.3333 17 14.6667 15.6667 16 15C14.6667 14.3333 13.3333 13 12 13Z" fill="var(--text-color)"></path>
            <path d="M12 10.5C11.1716 10.5 10.5 11.1716 10.5 12C10.5 12.8284 11.1716 13.5 12 13.5C12.8284 13.5 13.5 12.8284 13.5 12C13.5 11.1716 12.8284 10.5 12 10.5Z" fill="var(--primary-color)"></path>
        </svg>
    </div>
    """, unsafe_allow_html=True)
    
    # Campo de nombre opcional
    st.session_state.user_name = st.text_input("Tu Nombre (Opcional):", st.session_state.user_name)
    if st.session_state.user_name:
        st.markdown(f"¡Hola {st.session_state.user_name}!")
    
    st.markdown("---")
    
    # Botones de Navegación
    if st.button("1. Cargar Contenido", use_container_width=True):
        st.session_state.page = "Cargar Contenido"
    if st.button("2. Verificación IA", use_container_width=True):
        st.session_state.page = "Verificación IA"
    if st.button("3. Generar Examen", use_container_width=True):
        st.session_state.page = "Generar Examen"
    if st.button("4. Estudiar y Progreso", use_container_width=True):
        st.session_state.page = "Mi Progreso"
        
    st.markdown("---")
    
    # API Key de Gemini
    st.session_state.api_key = st.text_input("Google AI API Key", type="password", value=st.session_state.api_key)
    if st.session_state.api_key:
        st.success("API Key cargada.")
    else:
        st.info("Consigue tu API Key gratis en Google AI Studio.")

# --- CUERPO PRINCIPAL DE LA APP ---

# 1. Carga de Contenido
if st.session_state.page == "Cargar Contenido":
    st.header("1. Carga tu Contenido de Estudio 📚")
    st.markdown("Sube tus apuntes, resúmenes o presentaciones. Los analizaremos por ti.")
    
    uploaded_file = st.file_uploader(
        "Sube archivos .pdf, .pptx, .txt, .md",
        type=["pdf", "pptx", "txt", "md"],
        accept_multiple_files=False,
    )
    
    if uploaded_file:
        file_type = uploaded_file.type
        texto_extraido = ""
        
        with st.spinner(f"Procesando {uploaded_file.name}..."):
            try:
                if file_type == "application/pdf":
                    texto_extraido = extraer_texto_pdf(uploaded_file)
                elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    texto_extraido = extraer_texto_pptx(uploaded_file)
                elif file_type in ["text/plain", "text/markdown"]:
                    texto_extraido = uploaded_file.read().decode("utf-8")
                
                st.session_state.extracted_content = texto_extraido
                st.success("¡Archivo procesado con éxito!")
                st.info(f"Se extrajeron {len(texto_extraido)} caracteres.")
                
            except Exception as e:
                st.error(f"Ocurrió un error al procesar el archivo: {e}")
                st.session_state.extracted_content = None

    if st.session_state.extracted_content:
        st.subheader("Texto Extraído (Primeros 1000 caracteres):")
        st.text_area("", st.session_state.extracted_content[:1000] + "...", height=300)

# 2. Verificación Médica
elif st.session_state.page == "Verificación IA":
    st.header("2. Verificación Médica con IA 🔬")
    st.markdown("Analizamos la precisión científica de tu contenido.")

    if not st.session_state.extracted_content:
        st.warning("Por favor, carga un archivo primero en la pestaña 'Cargar Contenido'.")
    elif not st.session_state.api_key:
        st.warning("Por favor, introduce tu Google AI API Key en la barra lateral para continuar.")
    else:
        st.subheader("Contenido a Verificar:")
        st.text_area("", st.session_state.extracted_content, height=300, key="verif_content")
        
        if st.button("🔬 Analizar Precisión"):
            # --- CONEXIÓN REAL A GEMINI API ---
            try:
                genai.configure(api_key=st.session_state.api_key)
                model = genai.GenerativeModel(model_name="gemini-2.5-flash-preview-09-2025")
                
                prompt_parts = [
                    "Rol: Eres un profesor de medicina y revisor científico experto.",
                    f"Texto a revisar:\n---\n{st.session_state.extracted_content}\n---\n",
                    "Tu Tarea: Analiza el texto y evalúa su precisión científica, coherencia y claridad.",
                    "Marca los conceptos clave con un color/ícono:",
                    "🟢 Correcto y claro.",
                    "🟡 Parcialmente correcto (requiere aclaración).",
                    "🔴 Incorrecto o confuso.",
                    "Provee un resumen de tu análisis en formato Markdown.",
                    "Para puntos 🟡 y 🔴, provee una breve sugerencia o corrección con referencia a fuentes médicas estándar (ej. Harrison, ILAE, etc.)."
                ]

                with st.spinner("🧠 La IA está analizando la precisión..."):
                    response = model.generate_content(prompt_parts)
                    st.subheader("Resultados del Análisis de Gemini:")
                    st.markdown(response.text)

            except Exception as e:
                st.error(f"Error al conectar con Gemini: {e}")
                st.error("Asegúrate de que la API Key sea correcta.")

# 3. Generador de Preguntas (Página de CREACIÓN)
elif st.session_state.page == "Generar Examen":
    st.header("3. Generar Mazo de Flashcards 🎓")
    st.markdown("Crea un nuevo mazo de tarjetas de estudio basado en tu material.")

    if not st.session_state.extracted_content:
        st.warning("Por favor, carga un archivo primero para generar preguntas sobre él.")
    elif not st.session_state.api_key:
        st.warning("Por favor, introduce tu Google AI API Key en la barra lateral para continuar.")
    else:
        # Nuevo campo para el nombre del mazo
        deck_name = st.text_input("Nombre del Tema (ej. Fisiología Cardíaca - Ciclo):")
        
        st.markdown("---")
        
        col1, col2, col3 = st.columns(3)
        with col1:
            st.session_state.difficulty = st.selectbox("Nivel de Dificultad:", ["Automático (Adaptativo)", "Fácil", "Medio", "Difícil"])
        with col2:
            st.session_state.subject = st.selectbox("Tipo de Materia:", ["Materias Básicas (Anatomía, Fisio...)", "Materias Clínicas (Neuro, Pediatría...)"])
        with col3:
            st.session_state.num_questions = st.number_input("Número de Preguntas:", min_value=1, max_value=10, value=5)

        
        if st.button("🚀 Generar y Guardar Mazo"):
            # Validaciones
            if not deck_name:
                st.warning("Por favor, dale un nombre a tu mazo de tarjetas.")
            elif deck_name in st.session_state.flashcard_library:
                st.error(f"Ya existe un mazo con el nombre '{deck_name}'. Por favor, elige otro nombre.")
            else:
                # Limpiar el examen anterior
                restart_exam()
                
                # --- CONEXIÓN REAL A GEMINI API para MÚLTIPLES PREGUNTAS ---
                    f"Contexto del Estudiante: Nivel {st.session_state.difficulty}, Materia {st.session_state.subject}.",
                    f"Texto base (Material de estudio):\n---\n{st.session_state.extracted_content}\n---\n",
                    f"Tu Tarea: Genera {st.session_state.num_questions} preguntas de opción múltiple (4 opciones) basadas *únicamente* en el texto base.",
                    "Las preguntas deben ser claras, concisas y relevantes al estilo de examen médico.",
                    "Formato de Respuesta: Responde OBLIGATORIAMENTE en formato JSON. La estructura debe ser una LISTA de objetos:",
                    """
                    [
                      {
                        "pregunta": "El texto completo de la pregunta 1...",
                        "opciones": {
                          "A": "Texto de la opción A",
                          "B": "Texto de la opción B",
                          "C": "Texto de la opción C",
                          "D": "Texto de la opción D"
                        },
                        "respuesta_correcta": "B",
                        "explicacion": "Una breve pero completa explicación médica..."
                      },
                      {
                        "pregunta": "El texto completo de la pregunta 2...",
                        "opciones": { ... },
                        "respuesta_correcta": "A",
                        "explicacion": "..."
                      }
                    ]
                    """
                ]

                with st.spinner(f"🧠 Gemini está creando tu examen de {st.session_state.num_questions} preguntas..."):
                    response = model.generate_content(prompt_parts)
                    clean_response = response.text.strip().replace('```json', '').replace('```', '')
                    preguntas_json_list = json.loads(clean_response)
                    
                    # Guardar en la biblioteca en lugar de iniciar el examen
                    st.session_state.flashcard_library[deck_name] = preguntas_json_list
                    st.success(f"¡Mazo '{deck_name}' con {len(preguntas_json_list)} tarjetas guardado con éxito!")
                    st.balloons()

                except Exception as e:
                    st.error(f"Error al generar el examen con Gemini: {e}")
                    st.error("Asegúrate de que la API Key sea correcta y el modelo JSON haya funcionado.")
                    st.error(f"Respuesta recibida (para depuración): {response.text if 'response' in locals() else 'No response'}")

# --- PÁGINA DE ESTUDIO (NUEVA) ---
elif st.session_state.page == "Estudiar":
    
    if st.button("⬅️ Volver a mis mazos"):
        st.session_state.page = "Mi Progreso"
        restart_exam() # Limpia el estado del examen actual
        st.rerun()

    # --- Lógica para mostrar el examen (pregunta por pregunta) ---
    if st.session_state.current_exam:
        
        exam = st.session_state.current_exam
        idx = st.session_state.current_question_index
        
        # Verificar si el examen ha terminado
        if idx >= len(exam):
            st.header("¡Examen Completado! 🥳")
            
            correctas = sum(1 for r in st.session_state.exam_results if r['correcta'])
            total = len(exam)
            
            if total > 0:
                puntaje = (correctas / total) * 100
            else:
                puntaje = 0 # Evitar división por cero

            st.metric("Tu Puntaje:", f"{puntaje:.0f}%", f"{correctas} de {total} correctas")
            
            # Gráfico de pastel (pie chart) para el resumen final
            labels = ['Correctas', 'Incorrectas']
            values = [correctas, total - correctas]
            colors = ['#28a745', '#dc3545'] # Verde y Rojo

            fig = go.Figure(data=[go.Pie(labels=labels, values=values, hole=.3, 
                                        marker_colors=colors, 
                                        hoverinfo="label+percent+value",
                                        textinfo='percent',
                                        pull=[0, 0.05] # Separar un poco las incorrectas
                                        )])
            fig.update_layout(title_text='Resumen de Respuestas', title_x=0.5,
                              paper_bgcolor='rgba(0,0,0,0)',  # Fondo transparente
                              plot_bgcolor='rgba(0,0,0,0)',
                              font_color=var(--dark-text) # Color de texto
                              )
            st.plotly_chart(fig, use_container_width=True)

            st.subheader("Revisión Detallada:")
            for i, result in enumerate(st.session_state.exam_results):
                question_card = exam[i]
                if result['correcta']:
                    st.markdown(f"""
                    <div class="feedback-correct">
                        ✅ <strong>Pregunta {i+1} - ¡Correcto!</strong> (Seleccionaste: {result['seleccionada']})
                    </div>
                    """, unsafe_allow_html=True)
                else:
                    st.markdown(f"""
                    <div class="feedback-incorrect">
                        ❌ <strong>Pregunta {i+1} - Incorrecto.</strong> (Seleccionaste: {result['seleccionada']})
                        <br>
                        <strong>La respuesta correcta era:</strong> {result['correcta_texto']}
                    </div>
                    """, unsafe_allow_html=True)
                st.markdown(f"""
                <div class="feedback-explanation">
                    🧠 <strong>Explicación de la Pregunta {i+1}:</strong>
                    <br>
                    {question_card['explicacion']}
                </div>
                """, unsafe_allow_html=True)
                if st.button("Volver a mis mazos", on_click=restart_exam):
                    st.session_state.page = "Mi Progreso"
                    st.rerun() 
        
        else:
            # Mostrar la pregunta actual
            card = exam[idx]
            st.subheader(f"Tu Examen: Pregunta {idx + 1} de {len(exam)}")
            
            st.markdown('<div class="flashcard">', unsafe_allow_html=True)
            # La pregunta AHORA se renderiza correctamente dentro de la tarjeta
            st.markdown(f"<h5>{card['pregunta']}</h5>", unsafe_allow_html=True)
            
            opciones = list(card["opciones"].values())
            
            st.radio(
                "Selecciona tu respuesta:", 
                options=opciones,
                key=f"user_answer_{idx}", # Clave única por pregunta
                disabled=st.session_state.show_explanation
            )
            
            st.markdown('</div>', unsafe_allow_html=True) # Cierra la tarjeta
            
            # Botón de Responder (solo si no se ha respondido)
            if not st.session_state.show_explanation:
                if st.button("Responder y ver explicación"):
                    # Capturar la respuesta del radio button (usa la clave única)
                    current_user_selection = st.session_state[f"user_answer_{idx}"]
                    
                    if current_user_selection: # Asegurarse de que el usuario haya seleccionado algo
                        st.session_state.user_answer = current_user_selection # Actualizar el estado global con la selección actual
                        st.session_state.show_explanation = True
                        
                        user_ans_text = st.session_state.user_answer
                        correct_ans_letter = card["respuesta_correcta"]
                        correct_ans_text = card["opciones"][correct_ans_letter]
                        es_correcta = (user_ans_text == correct_ans_text)
                        
                        st.session_state.exam_results.append({
                            'correcta': es_correcta,
                            'seleccionada': user_ans_text,
                            'correcta_texto': correct_ans_text
                        })
                        
                        st.rerun() # Volver a cargar para mostrar la explicación
                    else:
                        st.warning("Por favor, selecciona una respuesta antes de continuar.")
            
            # Mostrar explicación y botón "Siguiente" (solo si ya se respondió)
            if st.session_state.show_explanation:
                result = st.session_state.exam_results[idx]
                if result['correcta']:
                    st.markdown(f"""
                    <div class="feedback-correct">
                        ✅ <strong>¡Correcto!</strong> La respuesta es: {result['correcta_texto']}
                    </div>
                    """, unsafe_allow_html=True)
                else:
                    st.markdown(f"""
                    <div class="feedback-incorrect">
                        ❌ <strong>Respuesta incorrecta.</strong> Seleccionaste: '{result['seleccionada']}'.
                        <br>
                        <strong>La respuesta correcta era:</strong> {result['correcta_texto']}
                    </div>
                    """, unsafe_allow_html=True)
                
                st.markdown(f"""
                <div class="feedback-explanation">
                    🧠 <strong>Explicación:</strong>
                    <br>
                    {card['explicacion']}
                </div>
                """, unsafe_allow_html=True)
                
                st.button("Siguiente Pregunta ➡️", on_click=go_to_next_question)

# 4. Progreso y Gamificación
elif st.session_state.page == "Mi Progreso":
    st.header("4. Estudiar y Progreso 🏆")
    
    st.subheader("Mis Mazos de Estudio 📚")
    
    # Lógica para seleccionar y empezar a estudiar un mazo
    if not st.session_state.flashcard_library:
        st.info("Aún no has generado ningún mazo. Ve a 'Generar Examen' para crear uno.")
    else:
        col1, col2 = st.columns([2, 1])
        with col1:
            deck_names = list(st.session_state.flashcard_library.keys())
            selected_deck_name = st.selectbox("Selecciona un mazo para estudiar:", options=deck_names)
        
        with col2:
            # Botón para iniciar el estudio
            if st.button("Iniciar Estudio 🚀", use_container_width=True, type="primary"):
                restart_exam() # Limpia el estado del examen anterior
                st.session_state.current_exam = st.session_state.flashcard_library[selected_deck_name]
                st.session_state.page = "Estudiar"
                st.rerun()

            # Botón para eliminar un mazo
            if st.button("🗑️ Eliminar Mazo", use_container_width=True):
                del st.session_state.flashcard_library[selected_deck_name]
                st.rerun()

    st.markdown("---") # Separador
    
    st.markdown("¡Sigue tu avance y colecciona insignias!")
    
    st.subheader("Niveles de Conocimiento")
    st.markdown("Básico ➔ Intermedio ➔ **Clínico** ➔ Experto")
    st.progress(75) # Ejemplo
    
    st.subheader("Mis Insignias")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown('<div class="doodle-container">🧬 Dominio Neuro</div>', unsafe_allow_html=True)
    with col2:
        st.markdown('<div class="doodle-container">❤️ Fisio Cardíaca</div>', unsafe_allow_html=True)
    with col3:
        st.markdown('<div class="doodle-container">🧪 Bioquímica</div>', unsafe_allow_html=True)

    st.subheader("Estadísticas de Desempeño")
    st.bar_chart({"Correctas": [20, 35, 30], "Incorrectas": [10, 5, 8]}, use_container_width=True)


